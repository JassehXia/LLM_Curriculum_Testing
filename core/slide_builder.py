import os
from typing import Union, cast
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.shapes.autoshape import Shape
from .schemas import SlideDeckSchema

def build_pptx_deck(deck_data: Union[SlideDeckSchema, dict], output_path: str) -> str:
    """
    Builds a PowerPoint presentation (.pptx) from a SlideDeckSchema object or dict.
    
    Args:
        deck_data: SlideDeckSchema object or dict adhering to SlideDeckSchema structure.
        output_path: Target path to save the .pptx file.
        
    Returns:
        The path to the generated presentation.
    """
    if isinstance(deck_data, dict):
        deck_schema = SlideDeckSchema(**deck_data)
    else:
        deck_schema = deck_data

    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]

    # Modern color palette
    NAVY = RGBColor(15, 23, 42)          # Slate 900
    WHITE = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(30, 41, 59)      # Slate 800
    CODE_BG = RGBColor(15, 23, 42)        # Slate 900 for code box background
    CODE_TEXT = RGBColor(226, 232, 240)    # Slate 200 for code text

    # ----------------------------------------------------
    # 1. Title Slide
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Title Slide Background Shape
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = NAVY
    bg_shape.line.color.rgb = NAVY

    # Title Text Frame
    tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf = tx_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = deck_schema.deck_title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.text = "Autonomous Curriculum Generator — PyTorch Learning Module"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
    p2.font.name = "Calibri"
    p2.space_before = Pt(20)

    # ----------------------------------------------------
    # 2. Content Slides
    # ----------------------------------------------------
    for slide_info in deck_schema.slides:
        slide = prs.slides.add_slide(blank_slide_layout)

        # Header Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_info.title
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY
        p_title.font.name = "Calibri"

        snippet = slide_info.code_snippet.strip() if slide_info.code_snippet else ""
        has_code = len(snippet) > 0
        content_width = Inches(5.6) if has_code else Inches(11.7)

        # Bullet Points Box
        bullets_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), content_width, Inches(5.0))
        tf_bullets = bullets_box.text_frame
        tf_bullets.word_wrap = True

        for i, bp in enumerate(slide_info.bullet_points):
            p_bp = tf_bullets.paragraphs[0] if i == 0 else tf_bullets.add_paragraph()
            p_bp.text = f"• {bp}"
            p_bp.font.size = Pt(18)
            p_bp.font.color.rgb = DARK_TEXT
            p_bp.font.name = "Calibri"
            p_bp.space_before = Pt(12)

        # Optional Code Snippet Panel
        if has_code:
            code_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0)
            )
            code_bg.fill.solid()
            code_bg.fill.fore_color.rgb = CODE_BG
            code_bg.line.color.rgb = CODE_BG

            code_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.9), Inches(5.5), Inches(4.8))
            tf_code = code_box.text_frame
            tf_code.word_wrap = True
            
            p_code = tf_code.paragraphs[0]
            p_code.text = snippet
            p_code.font.size = Pt(13)
            p_code.font.name = "Consolas"
            p_code.font.color.rgb = CODE_TEXT

    # Save presentation
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path
